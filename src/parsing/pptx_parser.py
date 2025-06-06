from pptx import Presentation
from pathlib import Path

def parse_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    result = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            result.append({
                "text": "\n".join(texts),
                "source": Path(path).name,
                "page": i  # слайды считаем как страницы
            })

    return result

#if __name__ == "__main__":
    from pprint import pprint

    result = parse_pptx("data/raw/test_1.pptx")  # замени на свой путь
    pprint(result[:2])  # покажет первые 2 чанка